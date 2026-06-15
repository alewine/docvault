"""Text extraction / OCR.

Per-format text extractors and the `ocr_file` dispatcher that routes an
uploaded file to the right one. Everything here is pure extraction: given a
file on disk, produce a readable text string. No DB writes (beyond the
`log_event` processing trace), no enrichment (title/summary/category/date
inference), no embedding/chunking.

Optional heavyweight deps (cv2, numpy, PIL, pytesseract, pdf2image,
pdfplumber, pillow_heif, python-docx, openpyxl, python-pptx) are imported
lazily inside the functions that need them, to keep process startup cheap.
"""

import sqlite3
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from config import TEXT_BASED_EXTENSIONS
from db import log_event


def preprocess_document_image(
    image_input,
    doc_id: str | None = None,
    conn: sqlite3.Connection | None = None,
):
    """
    Perspective-correct, upscale, and binarise a document image before Tesseract.
    Accepts a file path (str/Path) or a PIL Image. Returns a PIL Image.
    """
    import cv2
    import numpy as np
    from PIL import Image

    # Load into a BGR numpy array
    if isinstance(image_input, (str, Path)):
        bgr = cv2.imread(str(image_input))
        if bgr is None:
            raise ValueError(f"cv2.imread could not load: {image_input}")
    else:
        # PIL Image (e.g. from HEIC conversion) → BGR
        bgr = cv2.cvtColor(np.array(image_input.convert("RGB")), cv2.COLOR_RGB2BGR)

    h, w = bgr.shape[:2]

    # --- 1. Perspective correction ---
    gray = cv2.cvtColor(bgr, cv2.COLOR_BGR2GRAY)
    blurred = cv2.GaussianBlur(gray, (5, 5), 0)
    edges = cv2.Canny(blurred, 75, 200)
    contours, _ = cv2.findContours(edges, cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
    contours = sorted(contours, key=cv2.contourArea, reverse=True)

    quad = None
    for contour in contours[:5]:
        perimeter = cv2.arcLength(contour, True)
        approx = cv2.approxPolyDP(contour, 0.02 * perimeter, True)
        if len(approx) == 4 and cv2.contourArea(contour) > 0.20 * h * w:
            quad = approx
            break

    if quad is not None:
        corners = quad.reshape(4, 2).astype(np.float32)
        s = corners.sum(axis=1)
        d = np.diff(corners, axis=1).ravel()
        ordered = np.array([
            corners[np.argmin(s)],   # top-left
            corners[np.argmin(d)],   # top-right
            corners[np.argmax(s)],   # bottom-right
            corners[np.argmax(d)],   # bottom-left
        ], dtype=np.float32)
        tl, tr, br, bl = ordered
        out_w = int(max(np.linalg.norm(br - bl), np.linalg.norm(tr - tl)))
        out_h = int(max(np.linalg.norm(tr - br), np.linalg.norm(tl - bl)))
        dst = np.array(
            [[0, 0], [out_w - 1, 0], [out_w - 1, out_h - 1], [0, out_h - 1]],
            dtype=np.float32,
        )
        M = cv2.getPerspectiveTransform(ordered, dst)
        bgr = cv2.warpPerspective(bgr, M, (out_w, out_h))
        h, w = bgr.shape[:2]
        if doc_id and conn:
            log_event(conn, doc_id, "preprocess", status="success",
                      message=f"Perspective correction applied, corners: {ordered.tolist()}")
    else:
        if doc_id and conn:
            log_event(conn, doc_id, "preprocess", status="skipped",
                      message="No document boundary detected, skipping perspective correction")

    # --- 2. Upscaling ---
    if min(w, h) < 2400:
        scale = 2400 / min(w, h)
        new_w, new_h = int(w * scale), int(h * scale)
        bgr = cv2.resize(bgr, (new_w, new_h), interpolation=cv2.INTER_CUBIC)
        if doc_id and conn:
            log_event(conn, doc_id, "preprocess",
                      message=f"Upscaled from {w}x{h} to {new_w}x{new_h}")
        w, h = new_w, new_h

    # --- 3. Thresholding ---
    # Adaptive threshold works well on perspective-corrected crops; Otsu handles
    # full phone photos better because adaptive method destroys light-weight text
    # when there's no tight document boundary to normalize against.
    gray2 = cv2.cvtColor(bgr, cv2.COLOR_BGR2GRAY)
    if quad is not None:
        processed = cv2.adaptiveThreshold(
            gray2, 255,
            cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
            cv2.THRESH_BINARY,
            11, 2,
        )
        if doc_id and conn:
            log_event(conn, doc_id, "preprocess", message="Adaptive threshold applied (post-warp)")
    else:
        _, processed = cv2.threshold(gray2, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        if doc_id and conn:
            log_event(conn, doc_id, "preprocess", message="Otsu threshold applied (no boundary detected)")

    return Image.fromarray(processed)


def _best_psm_ocr(image_input) -> str:
    """Try PSM 6, 4, 3 and return the result with the most characters."""
    import pytesseract

    best = ""
    for psm in (6, 4, 3):
        try:
            text = pytesseract.image_to_string(image_input, config=f"--psm {psm}")
            if len(text.strip()) > len(best.strip()):
                best = text
        except Exception:
            pass
    return best


def _reconstruct_table_text(image_input) -> tuple[str | None, float | None]:
    """
    Use word-level bounding boxes to reconstruct tabular layout.
    Returns (None, None) when too few high-confidence tokens are found, signalling
    the caller to fall back to plain image_to_string().
    Returns (text, avg_confidence) on success.
    """
    import pytesseract
    from pytesseract import Output

    data = pytesseract.image_to_data(image_input, output_type=Output.DICT)

    words = []
    confidences: list[float] = []
    for i in range(len(data["text"])):
        conf = int(data["conf"][i])
        text = data["text"][i].strip()
        if conf < 40 or not text:
            continue
        words.append({
            "x": data["left"][i],
            "y": data["top"][i],
            "h": data["height"][i],
            "text": text,
        })
        confidences.append(conf)

    if len(words) < 5:
        return None, None

    avg_conf = round(sum(confidences) / len(confidences), 1) if confidences else None

    # Cluster into rows by vertical centre with 10 px tolerance
    rows: list[list[dict]] = []
    for word in sorted(words, key=lambda w: w["y"]):
        center = word["y"] + word["h"] / 2
        placed = False
        for row in rows:
            row_center = sum(w["y"] + w["h"] / 2 for w in row) / len(row)
            if abs(center - row_center) <= 10:
                row.append(word)
                placed = True
                break
        if not placed:
            rows.append([word])

    lines = []
    for row in rows:
        row.sort(key=lambda w: w["x"])
        lines.append("\t".join(w["text"] for w in row))

    return "\n".join(lines), avg_conf


_NON_DOCUMENT_TOKEN_THRESHOLD = 10
_NON_DOCUMENT_CONF_THRESHOLD = 50


def _probe_image_for_text(preprocessed_img) -> tuple[int, float]:
    """
    Cheap pre-OCR probe on a downscaled image. Returns (token_count, avg_confidence)
    where tokens are words with Tesseract confidence > _NON_DOCUMENT_CONF_THRESHOLD.
    Used to distinguish portrait photos from printed documents before full OCR.
    """
    import pytesseract
    from PIL import Image
    from pytesseract import Output

    img = preprocessed_img if isinstance(preprocessed_img, Image.Image) else Image.open(str(preprocessed_img))
    w, h = img.size
    longest = max(w, h)
    if longest > 800:
        scale = 800 / longest
        img = img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)

    data = pytesseract.image_to_data(img, output_type=Output.DICT)
    high_conf_tokens = [
        int(data["conf"][i])
        for i in range(len(data["text"]))
        if data["text"][i].strip() and int(data["conf"][i]) > _NON_DOCUMENT_CONF_THRESHOLD
    ]
    token_count = len(high_conf_tokens)
    avg_conf = round(sum(high_conf_tokens) / token_count, 1) if high_conf_tokens else 0.0
    return token_count, avg_conf


def _ocr_image(image_input, doc_id: str | None = None, conn: sqlite3.Connection | None = None, original_image=None) -> str:
    import pytesseract

    if doc_id and conn:
        log_event(conn, doc_id, "ocr_attempt", pipeline_path="tesseract_bbox",
                  message="Running Tesseract with bounding-box extraction")

    bbox_text, avg_conf = _reconstruct_table_text(image_input)
    if bbox_text is None:
        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="tesseract_bbox",
                      status="skipped", message="Insufficient high-confidence tokens, using plain OCR")
        text = _best_psm_ocr(image_input)
        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="tesseract_fallback",
                      status="success", char_count=len(text))
    else:
        text = bbox_text
        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="tesseract_bbox",
                      status="success", char_count=len(text),
                      message=f"Tesseract extracted {len(text)} chars",
                      metadata={"avg_confidence": avg_conf} if avg_conf is not None else None)

    if len(text) < 500 and original_image is not None:
        import cv2
        import numpy as np
        gray_original = cv2.cvtColor(np.array(original_image), cv2.COLOR_RGB2GRAY)
        fallback_text = pytesseract.image_to_string(gray_original, config='--psm 6')
        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="tesseract_gray_fallback",
                      char_count=len(fallback_text),
                      message=f"Gray fallback produced {len(fallback_text)} chars vs primary {len(text)} chars")
        if len(fallback_text) > len(text):
            text = fallback_text

    return text


_PDF_NATIVE_MIN_CHARS = 100


def _extract_pdf_native(path: Path, doc_id: str | None = None, conn: sqlite3.Connection | None = None) -> str | None:
    """
    Try to extract text from a PDF using pdfplumber (text layer + table layout).
    Returns None when the PDF appears to be a scanned image (< _PDF_NATIVE_MIN_CHARS
    of extractable text), so the caller can fall back to OCR.
    """
    import pdfplumber

    page_parts: list[str] = []
    total_tables = 0
    pages_with_tables: list[int] = []
    with pdfplumber.open(str(path)) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            parts: list[str] = []

            tables = page.extract_tables()
            if tables:
                total_tables += len(tables)
                pages_with_tables.append(page_num)
            for table in tables:
                rows = []
                for row in table:
                    rows.append("\t".join(cell or "" for cell in row))
                parts.append("\n".join(rows))

            prose = page.extract_text()
            if prose:
                parts.append(prose)

            if parts:
                page_parts.append("\n\n".join(parts))

    if doc_id and conn:
        if total_tables > 0:
            log_event(conn, doc_id, "table_extraction",
                      pipeline_path="pdfplumber",
                      message=f"Found {total_tables} table(s) across {len(pages_with_tables)} page(s)",
                      metadata={"table_count": total_tables, "pages_with_tables": pages_with_tables})
        else:
            log_event(conn, doc_id, "table_extraction",
                      pipeline_path="pdfplumber",
                      message="No tables detected")

    text = "\n\n".join(page_parts)
    if len(text.strip()) < _PDF_NATIVE_MIN_CHARS:
        return None
    return text


_JSON_MAX_CHARS = 200_000
_JSON_ARRAY_SAMPLE_HEAD = 50
_JSON_ARRAY_SAMPLE_TAIL = 10
_JSON_ARRAY_MAX = 500


def _flatten_json(obj: object, prefix: str, lines: list[str]) -> None:
    if isinstance(obj, dict):
        for k, v in obj.items():
            _flatten_json(v, f"{prefix}.{k}" if prefix else k, lines)
    elif isinstance(obj, list):
        for i, item in enumerate(obj):
            _flatten_json(item, f"{prefix}[{i}]", lines)
    elif obj is not None:
        lines.append(f"{prefix}: {obj}")


def extract_json_text(
    path: Path,
    doc_id: str | None = None,
    conn: sqlite3.Connection | None = None,
) -> str:
    import json as _json

    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        data = _json.loads(raw)
    except _json.JSONDecodeError as e:
        raise ValueError(f"Invalid JSON: {e}") from e

    truncation_note = ""
    if isinstance(data, list) and len(data) > _JSON_ARRAY_MAX:
        original_len = len(data)
        data = data[:_JSON_ARRAY_SAMPLE_HEAD] + data[-_JSON_ARRAY_SAMPLE_TAIL:]
        truncation_note = (
            f"[NOTE: Array truncated — showing first {_JSON_ARRAY_SAMPLE_HEAD} and "
            f"last {_JSON_ARRAY_SAMPLE_TAIL} of {original_len} items]\n\n"
        )

    lines: list[str] = []
    _flatten_json(data, "", lines)
    text = truncation_note + "\n".join(lines)

    if len(text) > _JSON_MAX_CHARS:
        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="json_parse",
                      status="warning",
                      message=f"JSON flattened output truncated at {_JSON_MAX_CHARS} chars")
        text = text[:_JSON_MAX_CHARS]

    return text


def ocr_file(path: Path, ext: str, doc_id: str | None = None, conn: sqlite3.Connection | None = None) -> str:
    from PIL import Image

    if ext == ".pdf":
        from pdf2image import convert_from_path

        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="pdfplumber",
                      message="Attempting native text extraction via pdfplumber")

        native_text = _extract_pdf_native(path, doc_id, conn)
        if native_text is None:
            if doc_id and conn:
                log_event(conn, doc_id, "ocr_attempt", pipeline_path="pdfplumber",
                          status="skipped", char_count=0,
                          message=f"pdfplumber returned text below threshold ({_PDF_NATIVE_MIN_CHARS} chars minimum), falling back to Tesseract")
            images = convert_from_path(str(path))
            if doc_id and conn:
                log_event(conn, doc_id, "ocr_attempt", pipeline_path="tesseract_bbox",
                          message=f"Running Tesseract with bounding-box extraction on {len(images)} page(s)")
            page_texts = []
            for img in images:
                preprocessed = preprocess_document_image(img, doc_id, conn)
                page_texts.append(_ocr_image(preprocessed, doc_id, conn, original_image=img))
            text = "\n\n".join(page_texts)
            if doc_id and conn:
                log_event(conn, doc_id, "ocr_attempt", pipeline_path="tesseract_fallback",
                          status="success", char_count=len(text),
                          message=f"OCR'd {len(images)} page(s)")
        else:
            text = native_text
            if doc_id and conn:
                log_event(conn, doc_id, "ocr_attempt", pipeline_path="pdfplumber",
                          status="success", char_count=len(text),
                          message=f"pdfplumber extracted {len(text)} chars")
    elif ext in (".heic", ".heif"):
        from pillow_heif import register_heif_opener
        register_heif_opener()
        try:
            img = Image.open(path).convert("RGB")
        except Exception as e:
            import traceback as _tb
            if doc_id and conn:
                log_event(conn, doc_id, "heic_conversion", status="failure",
                          message=f"{e}\n\n{_tb.format_exc()}")
            raise
        if doc_id and conn:
            log_event(conn, doc_id, "heic_conversion", status="success",
                      message="Converted HEIC to JPEG for OCR")
        preprocessed = preprocess_document_image(img, doc_id, conn)
        _tok, _conf = _probe_image_for_text(preprocessed)
        if _tok < _NON_DOCUMENT_TOKEN_THRESHOLD:
            if doc_id and conn:
                log_event(conn, doc_id, "non_document_image", status="skipped",
                          message=f"Non-document image: {_tok} high-confidence tokens (avg conf {_conf})",
                          metadata={"token_count": _tok, "avg_confidence": _conf})
            return ""
        text = _ocr_image(preprocessed, doc_id, conn, original_image=img)
    elif ext in (".jpg", ".jpeg", ".png"):
        orig_img = Image.open(str(path)).convert("RGB")
        preprocessed = preprocess_document_image(str(path), doc_id, conn)
        _tok, _conf = _probe_image_for_text(preprocessed)
        if _tok < _NON_DOCUMENT_TOKEN_THRESHOLD:
            if doc_id and conn:
                log_event(conn, doc_id, "non_document_image", status="skipped",
                          message=f"Non-document image: {_tok} high-confidence tokens (avg conf {_conf})",
                          metadata={"token_count": _tok, "avg_confidence": _conf})
            return ""
        text = _ocr_image(preprocessed, doc_id, conn, original_image=orig_img)
    elif ext == ".txt":
        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="txt_read",
                      message="Reading plain text file directly")
        text = path.read_text(encoding="utf-8", errors="replace")
    elif ext == ".csv":
        import csv as _csv
        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="csv_parse",
                      message="Parsing CSV file")
        rows = []
        with path.open(encoding="utf-8", errors="replace", newline="") as f:
            reader = _csv.reader(f)
            for row in reader:
                rows.append("\t".join(row))
        text = "\n".join(rows)
    elif ext == ".docx":
        import docx as _docx
        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="docx_parse",
                      message="Extracting text from DOCX")
        try:
            _doc = _docx.Document(str(path))
            parts = []
            for para in _doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            for table in _doc.tables:
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells]
                    parts.append("\t".join(row_cells))
            text = "\n".join(parts)
        except KeyError as exc:
            if "webSettings.xml" not in str(exc):
                raise
            if doc_id and conn:
                log_event(conn, doc_id, "ocr_attempt", pipeline_path="docx_parse",
                          status="skipped",
                          message=f"python-docx could not open optional DOCX part: {exc}; falling back to raw XML extraction")
            text = _extract_docx_text_from_archive(path)
    elif ext == ".xlsx":
        import openpyxl as _openpyxl
        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="xlsx_parse",
                      message="Extracting text from XLSX")
        wb = _openpyxl.load_workbook(str(path), data_only=True)
        sheet_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_parts.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(c.strip() for c in cells):
                    sheet_parts.append("\t".join(cells))
        text = "\n".join(sheet_parts)
    elif ext == ".pptx":
        from pptx import Presentation as _Presentation
        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="pptx_parse",
                      message="Extracting text from PPTX")
        prs = _Presentation(str(path))
        slide_parts = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_parts.append(t)
        text = "\n".join(slide_parts)
    elif ext == ".json":
        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="json_parse",
                      message="Parsing and flattening JSON file")
        text = extract_json_text(path, doc_id, conn)
        if doc_id and conn:
            log_event(conn, doc_id, "ocr_attempt", pipeline_path="json_parse",
                      status="success", char_count=len(text),
                      message=f"JSON flattened to {len(text)} chars")
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    text = text.strip()
    if ext in TEXT_BASED_EXTENSIONS and len(text) < 10:
        raise ValueError(
            f"Extracted text too short ({len(text)} chars) — file may be empty or unreadable"
        )
    return text


def _extract_docx_text_from_archive(path: Path) -> str:
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}

    with zipfile.ZipFile(path) as archive:
        document_xml = archive.read("word/document.xml")

    root = ET.fromstring(document_xml)
    parts: list[str] = []

    for para in root.findall(".//w:body//w:p", ns):
        texts = [node.text for node in para.findall(".//w:t", ns) if node.text]
        line = "".join(texts).strip()
        if line:
            parts.append(line)

    return "\n".join(parts)
